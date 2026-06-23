#!/usr/bin/env python3
"""Generuje firemní šablony PROKAT invest pro Word, Excel a PowerPoint."""

from __future__ import annotations

import io
import shutil
import zipfile
from copy import deepcopy
from pathlib import Path

from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Inches, Pt, RGBColor
from openpyxl import Workbook
from openpyxl.drawing.image import Image as XLImage
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.page import PageMargins
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt

BASE_DIR = Path(__file__).resolve().parent
ASSETS = BASE_DIR / "assets"
TEMPLATES = BASE_DIR / "templates"
CLIENT_DIR = BASE_DIR / "odeslat-klientovi"

LOGO = ASSETS / "logo-prokat-invest.png"
LOGO_SOURCE = ASSETS / "logo-prokat-invest-source.png"
WATERMARK = ASSETS / "watermark-prokat.png"

# Paleta z brand škály
NAVY = "10102A"
BRONZE = "4C3A2E"
GOLD_MID = "8E6B33"
GOLD = "D29C3A"
BLACK = "000000"
TEXT = "1A1A1A"
MUTED = "666666"
LIGHT_BG = "F5F5F5"
WHITE = "FFFFFF"

COMPANY = "PROKAT invest"
TAGLINE = "Investiční a projektové služby"
CONTACT_LINES = [
    "IČO: ____________________",
    "DIČ: ____________________",
    "E-mail: ____________________",
    "Telefon: ____________________",
    "Web: ____________________",
]


def hex_rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def prepare_logo(source: Path, target: Path) -> None:
  """Odstraní černé pozadí a uloží logo s průhledností."""
  img = Image.open(source).convert("RGBA")
  data = []
  for r, g, b, a in img.getdata():
    if r <= 4 and g <= 4 and b <= 4:
      data.append((r, g, b, 0))
    else:
      data.append((r, g, b, a))
  img.putdata(data)
  img.save(target, "PNG")


def make_watermark(source: Path, target: Path) -> None:
  """Vytvoří jemný vodoznak z loga."""
  img = Image.open(source).convert("RGBA")
  data = img.getdata()
  new_data = []
  for r, g, b, a in data:
    if a < 10:
      new_data.append((255, 255, 255, 0))
    else:
      gray = int(0.299 * r + 0.587 * g + 0.114 * b)
      new_data.append((gray, gray, gray, int(a * 0.14)))
  img.putdata(new_data)
  img.save(target, "PNG")


def set_cell_shading(cell, fill_hex: str) -> None:
  shading = OxmlElement("w:shd")
  shading.set(qn("w:fill"), fill_hex)
  shading.set(qn("w:val"), "clear")
  cell._tc.get_or_add_tcPr().append(shading)


def add_word_header_footer(doc: Document, title: str | None = None) -> None:
  section = doc.sections[0]
  section.top_margin = Cm(2.5)
  section.bottom_margin = Cm(2)
  section.left_margin = Cm(2.5)
  section.right_margin = Cm(2.5)

  header = section.header
  header.is_linked_to_previous = False
  hp = header.paragraphs[0]
  hp.clear()

  table = header.add_table(rows=1, cols=2, width=Inches(6.5))
  table.autofit = True
  left, right = table.rows[0].cells
  left.width = Inches(2.2)
  right.width = Inches(4.3)

  lp = left.paragraphs[0]
  lp.alignment = WD_ALIGN_PARAGRAPH.LEFT
  lp.add_run().add_picture(str(LOGO), width=Inches(2.4))

  rp = right.paragraphs[0]
  rp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
  name = rp.add_run(COMPANY)
  name.bold = True
  name.font.size = Pt(14)
  name.font.color.rgb = hex_rgb(BLACK)
  name.font.name = "Arial"

  if title:
    rp.add_run("\n")
    t = rp.add_run(title)
    t.font.size = Pt(10)
    t.font.color.rgb = hex_rgb(GOLD_MID)
    t.font.name = "Arial"

  accent_table = header.add_table(rows=1, cols=1, width=Inches(6.5))
  accent_cell = accent_table.rows[0].cells[0]
  accent_cell.text = ""
  set_cell_shading(accent_cell, GOLD)
  accent_cell.height = Pt(3)

  footer = section.footer
  footer.is_linked_to_previous = False
  fp = footer.paragraphs[0]
  fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
  fp.paragraph_format.space_before = Pt(6)
  text = fp.add_run(f"{COMPANY}  |  {TAGLINE}")
  text.font.size = Pt(8)
  text.font.color.rgb = hex_rgb(MUTED)
  text.font.name = "Arial"


def add_word_watermark(doc: Document) -> None:
  """Vloží jemný vodoznak loga do záhlaví."""
  header = doc.sections[0].header
  paragraph = header.add_paragraph()
  paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
  paragraph.paragraph_format.space_before = Pt(120)
  run = paragraph.add_run()
  run.add_picture(str(WATERMARK), width=Inches(4.5))


def add_word_styles(doc: Document) -> None:
  normal = doc.styles["Normal"]
  normal.font.name = "Arial"
  normal.font.size = Pt(11)
  normal.font.color.rgb = hex_rgb(TEXT)

  for level, size, color in [(1, 18, NAVY), (2, 14, BRONZE), (3, 12, GOLD_MID)]:
    style = doc.styles[f"Heading {level}"]
    style.font.name = "Arial"
    style.font.bold = True
    style.font.size = Pt(size)
    style.font.color.rgb = hex_rgb(color)


def create_word_document() -> None:
  doc = Document()
  add_word_styles(doc)
  add_word_header_footer(doc)
  add_word_watermark(doc)

  doc.add_paragraph()
  doc.add_heading("Nadpis dokumentu", level=1)
  doc.add_paragraph(
    "Toto je vzorový text. Sem pište obsah dopisu, smlouvy, zprávy nebo jiného dokumentu. "
    "Šablona obsahuje firemní záhlaví, patičku a vodoznak s logem PROKAT invest."
  )
  doc.add_heading("Podnadpis", level=2)
  doc.add_paragraph(
    "Odstavec s běžným formátováním. Pro úpravu stylů použijte styly Nadpis 1, Nadpis 2 a Odstavec."
  )

  doc.add_paragraph()
  doc.add_paragraph("S pozdravem,")
  doc.add_paragraph()
  sig = doc.add_paragraph("__________________________")
  sig.add_run("\nJméno a příjmení\nFunkce")

  out = TEMPLATES / "word" / "PROKAT-vzor-dokument.docx"
  doc.save(out)


def create_word_offer() -> None:
  doc = Document()
  add_word_styles(doc)
  add_word_header_footer(doc, title="Cenová nabídka")
  add_word_watermark(doc)

  p = doc.add_paragraph()
  p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
  run = p.add_run("Číslo nabídky: CN-____/2026\nDatum: __.__.2026\nPlatnost: 30 dní")
  run.font.size = Pt(10)
  run.font.color.rgb = hex_rgb(MUTED)

  doc.add_heading("Cenová nabídka", level=1)

  info = doc.add_table(rows=4, cols=2)
  info.style = "Table Grid"
  labels = ["Zákazník", "Adresa", "IČO / DIČ", "Kontaktní osoba"]
  for i, label in enumerate(labels):
    info.rows[i].cells[0].text = label
    info.rows[i].cells[1].text = ""
    for cell in info.rows[i].cells:
      for paragraph in cell.paragraphs:
        for r in paragraph.runs:
          r.font.name = "Arial"
          r.font.size = Pt(10)

  doc.add_paragraph()
  doc.add_heading("Položky nabídky", level=2)

  table = doc.add_table(rows=2, cols=5)
  table.style = "Table Grid"
  table.alignment = WD_TABLE_ALIGNMENT.CENTER
  headers = ["Položka", "Množství", "MJ", "Cena za MJ", "Celkem"]
  for idx, title in enumerate(headers):
    cell = table.rows[0].cells[idx]
    cell.text = title
    set_cell_shading(cell, NAVY)
    for paragraph in cell.paragraphs:
      paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
      for r in paragraph.runs:
        r.bold = True
        r.font.color.rgb = hex_rgb(WHITE)
        r.font.name = "Arial"
        r.font.size = Pt(10)

  sample = table.rows[1].cells
  sample[0].text = "Popis služby / produktu"
  sample[1].text = "1"
  sample[2].text = "ks"
  sample[3].text = "0,00 Kč"
  sample[4].text = "0,00 Kč"

  doc.add_paragraph()
  totals = doc.add_table(rows=3, cols=2)
  totals.alignment = WD_TABLE_ALIGNMENT.RIGHT
  for row, label in enumerate(["Celkem bez DPH:", "DPH 21 %:", "Celkem s DPH:"]):
    totals.rows[row].cells[0].text = label
    totals.rows[row].cells[1].text = "0,00 Kč"
    if row == 2:
      for cell in totals.rows[row].cells:
        set_cell_shading(cell, LIGHT_BG)
        for paragraph in cell.paragraphs:
          for r in paragraph.runs:
            r.bold = True

  doc.add_paragraph()
  doc.add_heading("Obchodní podmínky", level=2)
  doc.add_paragraph(
    "• Dodací lhůta: dle dohody\n"
    "• Platební podmínky: dle dohody\n"
    "• Nabídka je nezávazná do uplynutí platnosti"
  )

  doc.add_paragraph()
  doc.add_paragraph("Schvaluji nabídku:")
  doc.add_paragraph("Datum: ____________    Podpis: __________________________")

  out = TEMPLATES / "word" / "PROKAT-vzor-cenova-nabidka.docx"
  doc.save(out)


def thin_border(color: str = "CCCCCC") -> Border:
  side = Side(style="thin", color=color)
  return Border(left=side, right=side, top=side, bottom=side)


def style_header_row(ws, row: int, labels: list[str], start_col: int = 1) -> None:
  for idx, label in enumerate(labels):
    col = start_col + idx
    cell = ws.cell(row=row, column=col, value=label)
    cell.font = Font(name="Arial", bold=True, color=WHITE, size=10)
    cell.fill = PatternFill("solid", fgColor=NAVY)
    cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    cell.border = thin_border(NAVY)


def setup_excel_page(ws, title: str) -> None:
  ws.page_margins = PageMargins(left=0.6, right=0.6, top=0.75, bottom=0.75)
  ws.sheet_view.showGridLines = False
  ws.print_title_rows = "1:4"

  ws.merge_cells("B1:F1")
  ws.row_dimensions[1].height = 52
  logo = XLImage(str(LOGO))
  logo.width = 260
  logo.height = int(260 * logo.height / logo.width)
  ws.add_image(logo, "A1")

  ws["B1"] = COMPANY
  ws["B1"].font = Font(name="Arial", bold=True, size=16, color=BLACK)
  ws["B1"].alignment = Alignment(horizontal="right", vertical="center")

  ws.merge_cells("A2:F2")
  ws["A2"] = title
  ws["A2"].font = Font(name="Arial", size=11, color=GOLD_MID)
  ws["A2"].alignment = Alignment(horizontal="right")

  for col, width in zip("ABCDEF", [18, 14, 10, 14, 14, 18]):
    ws.column_dimensions[col].width = width

  ws.merge_cells("A3:F3")
  ws["A3"].fill = PatternFill("solid", fgColor=GOLD)
  ws.row_dimensions[3].height = 4

  watermark = XLImage(str(WATERMARK))
  watermark.width = 320
  watermark.height = int(320 * watermark.height / watermark.width)
  ws.add_image(watermark, "C8")


def create_excel_document() -> None:
  wb = Workbook()
  ws = wb.active
  ws.title = "Dokument"
  setup_excel_page(ws, TAGLINE)

  ws.merge_cells("A5:F5")
  ws["A5"] = "Nadpis dokumentu"
  ws["A5"].font = Font(name="Arial", bold=True, size=16, color=NAVY)

  ws.merge_cells("A7:F18")
  ws["A7"] = (
    "Sem pište obsah dokumentu. Tato šablona slouží jako univerzální firemní podklad "
    "pro texty, poznámky, zápisy nebo interní dokumenty společnosti PROKAT invest."
  )
  ws["A7"].font = Font(name="Arial", size=11, color=TEXT)
  ws["A7"].alignment = Alignment(wrap_text=True, vertical="top")

  ws.merge_cells("A20:F20")
  ws["A20"] = f"{COMPANY}  |  {TAGLINE}"
  ws["A20"].font = Font(name="Arial", size=8, color=MUTED)
  ws["A20"].alignment = Alignment(horizontal="center")

  out = TEMPLATES / "excel" / "PROKAT-vzor-dokument.xlsx"
  wb.save(out)


def create_excel_offer() -> None:
  wb = Workbook()
  ws = wb.active
  ws.title = "Cenová nabídka"
  setup_excel_page(ws, "Cenová nabídka")

  fields = [
    ("Číslo nabídky", "CN-____/2026"),
    ("Datum", "__.__.2026"),
    ("Platnost", "30 dní"),
    ("Zákazník", ""),
    ("Adresa", ""),
    ("IČO / DIČ", ""),
  ]
  row = 5
  for label, value in fields:
    ws.cell(row=row, column=1, value=label).font = Font(name="Arial", bold=True, size=10, color=BRONZE)
    ws.merge_cells(start_row=row, start_column=2, end_row=row, end_column=6)
    ws.cell(row=row, column=2, value=value).font = Font(name="Arial", size=10)
    ws.cell(row=row, column=2).border = thin_border()
    row += 1

  header_row = row + 1
  style_header_row(ws, header_row, ["Položka", "Množství", "MJ", "Cena za MJ", "Celkem"])

  data_start = header_row + 1
  for i in range(8):
    r = data_start + i
    ws.cell(r, 1, "Popis služby / produktu").border = thin_border()
    ws.cell(r, 2, 1).border = thin_border()
    ws.cell(r, 3, "ks").border = thin_border()
    ws.cell(r, 4, 0).number_format = '#,##0.00 "Kč"'
    ws.cell(r, 4).border = thin_border()
    ws.cell(r, 5, f"=B{r}*D{r}").number_format = '#,##0.00 "Kč"'
    ws.cell(r, 5).border = thin_border()

  total_row = data_start + 9
  ws.cell(total_row, 4, "Celkem bez DPH:").font = Font(bold=True)
  ws.cell(total_row, 5, f"=SUM(E{data_start}:E{data_start + 7})").number_format = '#,##0.00 "Kč"'
  ws.cell(total_row + 1, 4, "DPH 21 %:").font = Font(bold=True)
  ws.cell(total_row + 1, 5, f"=E{total_row}*0.21").number_format = '#,##0.00 "Kč"'
  ws.cell(total_row + 2, 4, "Celkem s DPH:").font = Font(bold=True, color=NAVY)
  ws.cell(total_row + 2, 5, f"=E{total_row}+E{total_row + 1}").number_format = '#,##0.00 "Kč"'
  for r in range(total_row, total_row + 3):
    ws.cell(r, 5).fill = PatternFill("solid", fgColor=LIGHT_BG)

  note_row = total_row + 5
  ws.merge_cells(start_row=note_row, start_column=1, end_row=note_row + 2, end_column=6)
  ws.cell(
    note_row,
    1,
    "Obchodní podmínky:\n• Dodací lhůta: dle dohody\n• Platební podmínky: dle dohody",
  )
  ws.cell(note_row, 1).alignment = Alignment(wrap_text=True, vertical="top")
  ws.cell(note_row, 1).font = Font(name="Arial", size=10, color=MUTED)

  out = TEMPLATES / "excel" / "PROKAT-vzor-cenova-nabidka.xlsx"
  wb.save(out)


def ppt_color(value: str) -> PptRGBColor:
  value = value.lstrip("#")
  return PptRGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def send_shape_to_back(shape) -> None:
  tree = shape._element.getparent()
  tree.remove(shape._element)
  tree.insert(2, shape._element)


def add_brand_to_slide(slide, subtitle: str) -> None:
  watermark = slide.shapes.add_picture(
    str(WATERMARK),
    PptInches(2.3),
    PptInches(1.4),
    width=PptInches(5.2),
  )
  send_shape_to_back(watermark)

  slide.shapes.add_picture(str(LOGO), PptInches(0.4), PptInches(0.25), width=PptInches(2.3))

  accent = slide.shapes.add_shape(1, PptInches(0), PptInches(1.05), PptInches(10), PptInches(0.06))
  accent.fill.solid()
  accent.fill.fore_color.rgb = ppt_color(GOLD)
  accent.line.fill.background()

  footer_box = slide.shapes.add_textbox(PptInches(0.4), PptInches(7.0), PptInches(9), PptInches(0.3))
  footer_tf = footer_box.text_frame
  footer_tf.text = f"{COMPANY}  |  {subtitle}"
  footer_tf.paragraphs[0].font.size = PptPt(9)
  footer_tf.paragraphs[0].font.color.rgb = ppt_color(MUTED)


def brand_slide_master(prs: Presentation, subtitle: str) -> None:
  prs._brand_subtitle = subtitle


def add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
  slide = prs.slides.add_slide(prs.slide_layouts[0])
  add_brand_to_slide(slide, getattr(prs, "_brand_subtitle", TAGLINE))
  slide.shapes.title.text = title
  if slide.placeholders[1]:
    slide.placeholders[1].text = subtitle
  slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = ppt_color(NAVY)
  slide.shapes.title.text_frame.paragraphs[0].font.bold = True


def add_content_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
  slide = prs.slides.add_slide(prs.slide_layouts[1])
  add_brand_to_slide(slide, getattr(prs, "_brand_subtitle", TAGLINE))
  slide.shapes.title.text = title
  slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = ppt_color(BRONZE)
  body = slide.placeholders[1].text_frame
  body.clear()
  for idx, bullet in enumerate(bullets):
    p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
    p.text = bullet
    p.level = 0
    p.font.size = PptPt(16)


def create_powerpoint_document() -> None:
  prs = Presentation()
  brand_slide_master(prs, TAGLINE)
  add_title_slide(prs, "Název dokumentu", "Podtitul / datum")
  add_content_slide(
    prs,
    "Obsah",
    [
      "Úvodní bod prezentace",
      "Druhý bod — text doplňte dle potřeby",
      "Závěr nebo shrnutí",
    ],
  )
  add_content_slide(
    prs,
    "Další sekce",
    ["Sem doplňte vlastní obsah", "Šablona obsahuje vodoznak a firemní barvy"],
  )
  out = TEMPLATES / "powerpoint" / "PROKAT-vzor-dokument.pptx"
  prs.save(out)


def create_powerpoint_offer() -> None:
  prs = Presentation()
  brand_slide_master(prs, "Cenová nabídka")
  add_title_slide(prs, "Cenová nabídka", "Číslo: CN-____/2026  |  Datum: __.__.2026")

  slide = prs.slides.add_slide(prs.slide_layouts[5])
  add_brand_to_slide(slide, getattr(prs, "_brand_subtitle", "Cenová nabídka"))
  shapes = slide.shapes
  title = shapes.title
  title.text = "Údaje zákazníka"
  rows, cols = 5, 2
  left, top, width, height = PptInches(0.8), PptInches(1.6), PptInches(8.5), PptInches(2.5)
  table = shapes.add_table(rows, cols, left, top, width, height).table
  labels = ["Zákazník", "Adresa", "IČO / DIČ", "Kontakt", "Platnost nabídky"]
  for i, label in enumerate(labels):
    table.cell(i, 0).text = label
    table.cell(i, 1).text = ""
    table.cell(i, 0).text_frame.paragraphs[0].font.bold = True

  slide = prs.slides.add_slide(prs.slide_layouts[5])
  add_brand_to_slide(slide, getattr(prs, "_brand_subtitle", "Cenová nabídka"))
  slide.shapes.title.text = "Položky nabídky"
  table = slide.shapes.add_table(4, 5, PptInches(0.5), PptInches(1.5), PptInches(9), PptInches(2.8)).table
  headers = ["Položka", "Množství", "MJ", "Cena za MJ", "Celkem"]
  for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = ppt_color(NAVY)
    cell.text_frame.paragraphs[0].font.color.rgb = ppt_color(WHITE)
    cell.text_frame.paragraphs[0].font.bold = True
  for r in range(1, 4):
    table.cell(r, 0).text = "Popis služby / produktu"
    table.cell(r, 1).text = "1"
    table.cell(r, 2).text = "ks"

  add_content_slide(
    prs,
    "Obchodní podmínky",
    [
      "Dodací lhůta: dle dohody",
      "Platební podmínky: dle dohody",
      "Nabídka platí 30 dní od data vystavení",
    ],
  )

  out = TEMPLATES / "powerpoint" / "PROKAT-vzor-cenova-nabidka.pptx"
  prs.save(out)


TEMPLATE_CONTENT_TYPES = {
  ".docx": (
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.template.main+xml",
  ),
  ".xlsx": (
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.template.main+xml",
  ),
  ".pptx": (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
    "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
  ),
}

TEMPLATE_EXTENSIONS = {
  ".docx": ".dotx",
  ".xlsx": ".xltx",
  ".pptx": ".potx",
}


def convert_to_office_template(source: Path, target: Path) -> None:
  """Převede .docx/.xlsx/.pptx na oficiální formát šablony (.dotx/.xltx/.potx)."""
  suffix = source.suffix.lower()
  doc_type, template_type = TEMPLATE_CONTENT_TYPES[suffix]

  with zipfile.ZipFile(source, "r") as zin:
    contents = {name: zin.read(name) for name in zin.namelist()}

  content_types = contents["[Content_Types].xml"].decode("utf-8")
  content_types = content_types.replace(doc_type, template_type)
  contents["[Content_Types].xml"] = content_types.encode("utf-8")

  target.parent.mkdir(parents=True, exist_ok=True)
  with zipfile.ZipFile(target, "w", zipfile.ZIP_DEFLATED) as zout:
    for name, data in contents.items():
      zout.writestr(name, data)


def export_for_client() -> Path:
  """Zkopíruje šablony do složky pro klienta + vytvoří ZIP."""
  mapping = [
    (TEMPLATES / "word", ".docx", ".dotx"),
    (TEMPLATES / "excel", ".xlsx", ".xltx"),
    (TEMPLATES / "powerpoint", ".pptx", ".potx"),
  ]

  if CLIENT_DIR.exists():
    shutil.rmtree(CLIENT_DIR)
  CLIENT_DIR.mkdir(parents=True)

  packaged_files: list[Path] = []

  for source_dir, doc_ext, template_ext in mapping:
    client_subdir = CLIENT_DIR / source_dir.name
    client_subdir.mkdir(parents=True, exist_ok=True)

    for source in sorted(source_dir.glob(f"*{doc_ext}")):
      # Běžný Office soubor (.docx / .xlsx / .pptx) — otevře se přímo k úpravám
      doc_copy = client_subdir / source.name
      shutil.copy2(source, doc_copy)
      packaged_files.append(doc_copy)

      # Oficiální šablona (.dotx / .xltx / .potx) — při otevření vytvoří nový dokument
      template_name = source.stem + template_ext
      template_path = client_subdir / template_name
      convert_to_office_template(source, template_path)
      packaged_files.append(template_path)

  readme = CLIENT_DIR / "JAK-POUZIT-SABLONY.txt"
  readme.write_text(
    "PROKAT invest — firemní šablony\n"
    "================================\n\n"
    "Ve složkách word / excel / powerpoint najdete dva typy souborů:\n\n"
    "1) Soubory .docx / .xlsx / .pptx\n"
    "   Otevřou se přímo v Microsoft Office (nebo LibreOffice).\n"
    "   Vhodné pro okamžitou úpravu a odeslání klientovi.\n\n"
    "2) Soubory .dotx / .xltx / .potx  (oficiální šablony Office)\n"
    "   Po dvojkliku se vytvoří NOVÝ dokument podle vzoru.\n"
    "   Vhodné pro opakované použití ve firmě — vzor zůstane nedotčený.\n\n"
    "Doporučení:\n"
    "- Word šablony: Soubor → Uložit jako → Word šablona (.dotx)\n"
    "  nebo použijte dodané soubory .dotx\n"
    "- Excel: použijte .xltx pro nové nabídky, .xlsx pro jednorázovou práci\n"
    "- PowerPoint: použijte .potx pro nové prezentace\n\n"
    "Obsah složek:\n"
    "- PROKAT-vzor-dokument.*       → univerzální firemní dokument\n"
    "- PROKAT-vzor-cenova-nabidka.* → cenová nabídka s tabulkou položek\n",
    encoding="utf-8",
  )
  packaged_files.append(readme)

  zip_path = CLIENT_DIR / "PROKAT-Invest-sablony.zip"
  with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
    for path in packaged_files:
      zf.write(path, arcname=path.relative_to(CLIENT_DIR))

  return zip_path


def main() -> None:
  for folder in ("word", "excel", "powerpoint"):
    (TEMPLATES / folder).mkdir(parents=True, exist_ok=True)

  if not LOGO_SOURCE.exists():
    raise FileNotFoundError(f"Chybí zdrojové logo: {LOGO_SOURCE}")

  prepare_logo(LOGO_SOURCE, LOGO)
  make_watermark(LOGO, WATERMARK)

  create_word_document()
  create_word_offer()
  create_excel_document()
  create_excel_offer()
  create_powerpoint_document()
  create_powerpoint_offer()

  zip_path = export_for_client()

  print("Hotovo.")
  print("  Vzorové soubory:  ", TEMPLATES)
  print("  Pro odeslání:     ", CLIENT_DIR)
  print("  ZIP archiv:       ", zip_path)


if __name__ == "__main__":
  main()
